"""
kb_store.py
Turn a folder of requirement documents into a deterministic set of text
chunks for the Recursive Language Model to navigate. No embedding model
is involved: chunks are split on Markdown headings and paragraph
boundaries, given stable ids, and cached to kb_index.json keyed on file
mtimes so a rebuild only happens when a document actually changes.

Supported document types:
    .md .markdown .txt          read as text
    .pdf                        pypdf text extraction (OCR fallback for
                                scanned PDFs is applied by the indexer via
                                kb_ocr when an OCR engine is installed)
    .docx                       python-docx paragraphs + tables + headers
    .xlsx .xlsm                 openpyxl: every sheet, every non-empty row
    .pptx                       python-pptx: slide text + tables + notes
    .vsdx                       visio_convert: shapes, text, connections
    .html .htm                  selectolax text (regex fallback)
    .csv .tsv .json .jsonl .yaml .yml .xml .log   read as text

    MULTIMEDIA (optional backends; graceful degradation):
    .png .jpg .jpeg .gif .bmp .tiff .tif .webp .svg .ico .heic .heif
                                OCR for text in images + EXIF metadata
    .mp3 .wav .ogg .flac .m4a .wma .aac .opus .aiff
                                Speech-to-text (faster-whisper / whisper)
    .mp4 .mkv .avi .mov .webm .wmv .flv .m4v .mpg .mpeg .3gp
                                Audio transcription + keyframe OCR

    LEGACY / EXTENDED DOCUMENTS:
    .doc                        Word 97-2003 (antiword / raw OLE2 extraction)
    .ppt                        PowerPoint 97-2003 (OLE2 text records)
    .msg                        Outlook email (OLE2 structured storage)
    .odt .ods .odp .odg         OpenDocument (ZIP + XML content extraction)
    .eml .mbox                  Email messages (stdlib email parser)
    .tex .latex                 LaTeX source (command stripping)
    .epub .fb2                  Ebooks (XHTML/XML text extraction)
    .pages .numbers .key        Apple iWork (ZIP package text extraction)
    .wps .hwp                   WPS/Hangul (OLE2 + raw text runs)
    .man .info                  Unix man/info pages (troff stripping)

    anything else               read as text, but binary-looking content
                                (NUL bytes / many replacement chars) is
                                discarded rather than indexed as garbage

Public API:
    extract_text(path) -> str
    build_index(kb_dir) -> KbIndex
    load_or_build_index(kb_dir, index_path) -> KbIndex
    KbIndex.map_listing() -> str        compact id/title/size listing
    KbIndex.total_chars / .total_tokens
"""

from __future__ import annotations

import gc
import json
import re
import time
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Final

from core.app_config import RLM_MAP_CHUNK_TOKENS

# ~4 characters per token is the standard rough estimate.
_CHARS_PER_TOKEN: Final[int] = 4

# Bump whenever extraction logic changes in a way that should re-index
# existing files even if their bytes are unchanged. v2 added xlsx/pptx
# extraction, docx headers/footers, OCR-aware scanned-PDF handling, and a
# binary-garbage guard. v3 added Visio .vsdx diagram extraction.
# v4 added multimedia (image OCR, audio STT, video transcription).
# v5 added legacy document formats (.doc, .ppt, .odt, .eml, .epub, etc.)
EXTRACTOR_VERSION: Final[int] = 5

_TEXT_EXT: Final[frozenset[str]] = frozenset({
    ".md", ".markdown", ".txt", ".csv", ".tsv", ".json", ".jsonl",
    ".ndjson", ".yaml", ".yml", ".xml", ".log", ".rst", ".text",
})
# Archive formats whose contained files are extracted and indexed.
# .zip is implemented; the others are listed for future extension.
_ARCHIVE_EXT: Final[frozenset[str]] = frozenset(
    {".zip", ".7z", ".tar", ".gz", ".tgz"}
)
# Per-file text cap: prevents OOM on huge logs/dumps and bounds archive
# decompression (bomb guard). ~16 MB of text is far more than any doc needs.
_MAX_TEXT_BYTES: Final[int] = 16 * 1024 * 1024
_HEADING_RE: Final[re.Pattern[str]] = re.compile(r"^(#{1,6})\s+(.*)$")
_HTML_TAG_RE: Final[re.Pattern[str]] = re.compile(r"<[^>]+>")
_MULTI_BLANK_RE: Final[re.Pattern[str]] = re.compile(r"\n\s*\n\s*\n+")


def approx_tokens(text: str) -> int:
    return max(1, len(text) // _CHARS_PER_TOKEN)


# ---------------------------------------------------------------------
# Extraction
# ---------------------------------------------------------------------
def _extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        parts: list[str] = []
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                parts.append("")
        return "\n\n".join(parts)
    except Exception:
        return ""


def _extract_docx(path: Path) -> str:
    try:
        import docx  # python-docx
        doc = docx.Document(str(path))
        parts: list[str] = [p.text for p in doc.paragraphs if p.text]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                line = " | ".join(c for c in cells if c)
                if line:
                    parts.append(line)
        # Headers/footers often carry doc IDs, revision, and section context.
        try:
            for section in doc.sections:
                for hf in (section.header, section.footer):
                    for p in getattr(hf, "paragraphs", []):
                        if p.text and p.text.strip():
                            parts.append(p.text.strip())
        except Exception:
            pass
        return "\n".join(parts)
    except Exception:
        return ""


def _extract_xlsx(path: Path) -> str:
    """All sheets, all non-empty rows, as readable text. Read-only +
    data_only so large workbooks stream and formulas resolve to values."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        parts: list[str] = []
        for ws in wb.worksheets:
            parts.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [
                    ("" if v is None else str(v)).strip() for v in row
                ]
                if any(cells):
                    parts.append(" | ".join(c for c in cells if c))
        try:
            wb.close()
        except Exception:
            pass
        return "\n".join(parts)
    except Exception:
        return ""


def _extract_xls(path: Path) -> str:
    """Excel 97-2003 (.xls) via xlrd."""
    try:
        import xlrd
        wb = xlrd.open_workbook(str(path))
        parts: list[str] = []
        for ws in wb.sheets():
            parts.append(f"# Sheet: {ws.name}")
            for row_idx in range(ws.nrows):
                cells = [
                    str(ws.cell_value(row_idx, c)).strip()
                    for c in range(ws.ncols)
                ]
                if any(cells):
                    parts.append(" | ".join(c for c in cells if c))
        return "\n".join(parts)
    except Exception:
        return ""


def _extract_rtf(path: Path) -> str:
    """Rich Text Format via striprtf."""
    try:
        from striprtf.striprtf import rtf_to_text
        raw = path.read_bytes().decode("utf-8", errors="replace")
        return rtf_to_text(raw) or ""
    except Exception:
        return ""


def _extract_pptx(path: Path) -> str:
    """All slides: shape text, table cells, and speaker notes."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"# Slide {i}")
            for shape in slide.shapes:
                try:
                    if shape.has_text_frame:
                        t = shape.text_frame.text.strip()
                        if t:
                            parts.append(t)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [c.text.strip() for c in row.cells]
                            line = " | ".join(c for c in cells if c)
                            if line:
                                parts.append(line)
                except Exception:
                    continue
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"[Notes] {notes}")
            except Exception:
                pass
        return "\n".join(parts)
    except Exception:
        return ""


def _extract_visio(path: Path) -> str:
    """Visio .vsdx: shapes, text content, and connection structure."""
    try:
        from tools.visio_convert import extract_visio_text
        return extract_visio_text(path)
    except Exception:
        return ""


def _looks_binary(text: str) -> bool:
    """Heuristic: a sample with NUL bytes or a high share of U+FFFD
    replacement chars is binary read as text (e.g. an unparsed office/zip
    file). Used to avoid polluting the index with garbage."""
    if not text:
        return False
    sample = text[:4096]
    if "\x00" in sample:
        return True
    repl = sample.count("\ufffd")
    return repl > max(8, len(sample) // 20)


def _extract_html(path: Path) -> str:
    raw = _read_text_file(path)
    try:
        from selectolax.parser import HTMLParser
        return HTMLParser(raw).text(separator="\n").strip()
    except Exception:
        return _HTML_TAG_RE.sub(" ", raw)


def _read_text_file(path: Path) -> str:
    """Read text with a size cap to prevent OOM on huge logs/dumps."""
    try:
        size = path.stat().st_size
        if size > _MAX_TEXT_BYTES:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                return f.read(_MAX_TEXT_BYTES)
        return path.read_text(encoding="utf-8", errors="replace")
    except Exception:
        try:
            raw = path.read_bytes()[:_MAX_TEXT_BYTES]
            return raw.decode("utf-8", errors="replace")
        except Exception:
            return ""


def _extract_archive(path: Path) -> str:
    """Extract text from a .zip archive by iterating contained files,
    extracting each to a temp dir, and calling extract_text (skipping nested
    archives to prevent recursion bombs).
    Guards: path traversal (zip-slip), decompression bomb, total text cap."""
    import tempfile
    import zipfile as _zf

    if not _zf.is_zipfile(str(path)):
        return ""
    parts: list[str] = []
    total_chars = 0
    try:
        with _zf.ZipFile(str(path), "r") as zf:
            names = [n for n in zf.namelist()
                     if not n.endswith("/") and not n.startswith("__MACOSX")]
            with tempfile.TemporaryDirectory(prefix="kb_zip_") as td:
                td_path = Path(td)
                td_resolved = td_path.resolve()
                for name in names[:500]:  # cap at 500 files per archive
                    try:
                        member_ext = Path(name).suffix.lower()
                        # Skip nested archives to prevent recursion bombs
                        if member_ext in _ARCHIVE_EXT:
                            continue
                        # Guard: path traversal (zip-slip)
                        target = (td_path / name).resolve()
                        if not str(target).startswith(str(td_resolved)):
                            continue
                        # Guard: decompression bomb per member
                        info = zf.getinfo(name)
                        if info.file_size > _MAX_TEXT_BYTES:
                            continue
                        zf.extract(name, td)
                        member = td_path / name
                        if member.is_file() and member.stat().st_size > 0:
                            text = extract_text(member)
                            if text.strip():
                                parts.append(f"--- {name} ---\n{text}")
                                total_chars += len(text)
                                if total_chars > _MAX_TEXT_BYTES:
                                    break
                    except Exception:
                        continue
    except Exception:
        return ""
    result = "\n\n".join(parts)
    del parts
    gc.collect()
    return result


def extract_text(path: Path) -> str:
    """Best-effort plain-text extraction. Never raises; returns '' on
    failure so one bad file does not abort an index build."""
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext in (".xlsx", ".xlsm"):
        return _extract_xlsx(path)
    if ext == ".xls":
        return _extract_xls(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".vsdx":
        return _extract_visio(path)
    if ext == ".rtf":
        return _extract_rtf(path)
    if ext in (".html", ".htm"):
        return _extract_html(path)
    if ext == ".zip":
        return _extract_archive(path)
    if ext in _TEXT_EXT:
        return _read_text_file(path)
    # Multimedia formats: delegate to the isolated multimedia extractor.
    try:
        from kb.multimedia import MULTIMEDIA_EXTENSIONS, extract_multimedia_text
        if ext in MULTIMEDIA_EXTENSIONS:
            return extract_multimedia_text(path)
    except Exception:
        pass
    # Legacy/extended document formats (.doc, .ppt, .odt, .eml, .epub, etc.)
    try:
        from kb.legacy_docs import LEGACY_EXTENSIONS, extract_legacy_text
        if ext in LEGACY_EXTENSIONS:
            return extract_legacy_text(path)
    except Exception:
        pass
    # Unknown extension: try as text, but discard binary-looking content
    # (e.g. an unparsed .xls/.ppt/.zip/image) rather than indexing garbage.
    raw = _read_text_file(path)
    return "" if _looks_binary(raw) else raw


# ---------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------
@dataclass(slots=True)
class KbChunk:
    chunk_id: str
    doc: str
    title: str
    text: str
    n_chars: int = 0
    context: str = ""

    def __post_init__(self) -> None:
        if not self.n_chars:
            self.n_chars = len(self.text)

    @property
    def contextualized_text(self) -> str:
        """Text with the situating context prefix (if available)."""
        if self.context:
            return f"{self.context}\n\n{self.text}"
        return self.text


def _section_title(raw_title: str, body: str) -> str:
    if raw_title:
        return raw_title[:90]
    for line in body.splitlines():
        s = line.strip()
        if s:
            return s[:90]
    return "(untitled)"


def _split_sections(text: str) -> list[tuple[str, str]]:
    """Split on Markdown headings. Returns (title, body) pairs in order."""
    sections: list[tuple[str, str]] = []
    cur_title = ""
    cur_lines: list[str] = []
    for line in text.split("\n"):
        m = _HEADING_RE.match(line)
        if m:
            if cur_lines or cur_title:
                sections.append((cur_title, "\n".join(cur_lines).strip()))
            cur_title = m.group(2).strip()
            cur_lines = []
        else:
            cur_lines.append(line)
    if cur_lines or cur_title:
        sections.append((cur_title, "\n".join(cur_lines).strip()))
    return sections or [("", text.strip())]


def _split_by_budget(body: str, budget_chars: int) -> list[str]:
    """Accumulate paragraphs into chunks under budget_chars. Splits an
    oversized single paragraph on sentence-ish boundaries as a fallback."""
    if len(body) <= budget_chars:
        return [body] if body.strip() else []
    paras = [p for p in _MULTI_BLANK_RE.sub("\n\n", body).split("\n\n")]
    out: list[str] = []
    cur: list[str] = []
    cur_len = 0
    for para in paras:
        plen = len(para) + 2
        if plen > budget_chars and not cur:
            # Hard-split a giant paragraph.
            for i in range(0, len(para), budget_chars):
                out.append(para[i:i + budget_chars])
            continue
        if cur_len + plen > budget_chars and cur:
            out.append("\n\n".join(cur).strip())
            cur = [para]
            cur_len = plen
        else:
            cur.append(para)
            cur_len += plen
    if cur:
        out.append("\n\n".join(cur).strip())
    return [c for c in out if c.strip()]


def chunk_document(doc_name: str, doc_index: int, text: str) -> list[KbChunk]:
    budget = RLM_MAP_CHUNK_TOKENS * _CHARS_PER_TOKEN
    chunks: list[KbChunk] = []
    seq = 0
    for raw_title, body in _split_sections(text):
        pieces = _split_by_budget(body, budget)
        if not pieces and raw_title:
            pieces = [""]  # heading with no body still indexed
        for piece in pieces:
            title = _section_title(raw_title, piece)
            chunk_id = f"d{doc_index:03d}c{seq:04d}"
            full = (f"{raw_title}\n{piece}".strip()
                    if raw_title else piece.strip())
            chunks.append(KbChunk(
                chunk_id=chunk_id, doc=doc_name, title=title, text=full,
            ))
            seq += 1
    return chunks


# ---------------------------------------------------------------------
# Index
# ---------------------------------------------------------------------
@dataclass(slots=True)
class KbSource:
    name: str
    mtime: float
    size: int
    # Content SHA (short hex). Empty on indexes built before content-hash
    # incremental indexing; the currency check treats a missing/empty sha as
    # "needs rehash", which triggers exactly one rebuild on upgrade.
    sha: str = ""


@dataclass(slots=True)
class KbIndex:
    chunks: list[KbChunk] = field(default_factory=list)
    sources: list[KbSource] = field(default_factory=list)
    built_at: float = 0.0
    # Lazily-built chunk_id -> chunk map for O(1) lookups; not serialized.
    _id_cache: dict[str, KbChunk] | None = field(
        default=None, init=False, repr=False, compare=False
    )

    @property
    def total_chars(self) -> int:
        return sum(c.n_chars for c in self.chunks)

    @property
    def total_tokens(self) -> int:
        return approx_tokens("".join(c.text for c in self.chunks)) \
            if self.chunks else 0

    @property
    def n_docs(self) -> int:
        return len(self.sources)

    def by_id(self, chunk_id: str) -> KbChunk | None:
        """O(1) chunk lookup. The id->chunk map is built once on first use
        and reused; replaces the previous linear scan, which was O(n) per
        call and O(n*k) across the navigator's k selected ids."""
        cache = self._id_cache
        if cache is None:
            cache = {c.chunk_id: c for c in self.chunks}
            self._id_cache = cache
        return cache.get(chunk_id)

    def map_listing(self) -> str:
        """Compact navigable listing the navigator model reads to choose
        relevant chunks without seeing full bodies."""
        lines: list[str] = []
        last_doc = ""
        for c in self.chunks:
            if c.doc != last_doc:
                lines.append(f"== {c.doc} ==")
                last_doc = c.doc
            lines.append(
                f"  [{c.chunk_id}] (~{approx_tokens(c.text)} tok) {c.title}"
            )
        return "\n".join(lines)

    def to_dict(self) -> dict:
        return {
            "extractor_version": EXTRACTOR_VERSION,
            "built_at": self.built_at,
            "sources": [asdict(s) for s in self.sources],
            "chunks": [asdict(c) for c in self.chunks],
        }


# Prefer the structured spreadsheet over a same-content PDF "twin": for a
# config workbook present as both .xlsx and .pdf, the .xlsx extracts every
# sheet/row cleanly while pypdf flattens tables. Matching is on the core
# name after a leading "Category - " prefix (e.g. "Enhancement - ", "Fresh
# Update - "), so "Enhancement - CLEAN ... v8.2.xlsx" dedups the
# "Fresh Update - CLEAN ... v8.2.pdf". Only PDF-vs-spreadsheet pairs are
# touched; nothing else is ever dropped.
_SHEET_EXT: Final[frozenset[str]] = frozenset({".xlsx", ".xlsm"})


def _twin_key(p: Path) -> str:
    """Normalized core name for twin matching: stem minus a leading
    'Word(s) - ' category prefix, lowercased and whitespace-collapsed."""
    stem = p.stem
    if " - " in stem:
        stem = stem.split(" - ", 1)[1]
    return " ".join(stem.lower().split())


def dedup_twins(
    files: list[Path],
) -> tuple[list[Path], list[tuple[Path, Path]]]:
    """Return (kept, dropped_pairs). For each core name that has BOTH a
    spreadsheet (.xlsx/.xlsm) and one or more .pdf files, drop the PDF(s) and
    keep the spreadsheet. dropped_pairs is [(dropped_pdf, kept_sheet), ...]
    for logging. Deterministic; preserves input order for kept files."""
    by_key: dict[str, dict[str, list[Path]]] = {}
    for p in files:
        ext = p.suffix.lower()
        bucket = by_key.setdefault(_twin_key(p), {})
        bucket.setdefault(ext, []).append(p)
    drop: set[Path] = set()
    dropped_pairs: list[tuple[Path, Path]] = []
    for bucket in by_key.values():
        sheets = [q for e in _SHEET_EXT for q in bucket.get(e, [])]
        pdfs = bucket.get(".pdf", [])
        if sheets and pdfs:
            keep_sheet = sheets[0]
            for pdf in pdfs:
                drop.add(pdf)
                dropped_pairs.append((pdf, keep_sheet))
    kept = [p for p in files if p not in drop]
    return kept, dropped_pairs


def _raw_scan(kb_dir: Path) -> list[Path]:
    if not kb_dir.exists():
        return []
    return [
        p for p in sorted(kb_dir.rglob("*"), key=lambda x: str(x).lower())
        if p.is_file() and not p.name.startswith(".")
    ]


def _scan_sources(kb_dir: Path) -> list[Path]:
    """Indexable files for a KB dir, with PDF/spreadsheet twins deduped
    (spreadsheet preferred). Single source of truth so the index, the
    currency check, and counts all agree."""
    return dedup_twins(_raw_scan(kb_dir))[0]


def _source_records(
    files: list[Path], cache: dict | None = None,
) -> list[KbSource]:
    from kb.file_sig import file_sha

    cache = cache if cache is not None else {"entries": {}}
    out: list[KbSource] = []
    for p in files:
        try:
            st = p.stat()
            out.append(KbSource(
                name=p.name, mtime=st.st_mtime, size=st.st_size,
                sha=file_sha(p, cache),
            ))
        except OSError:
            out.append(KbSource(name=p.name, mtime=0.0, size=0, sha=""))
    return out


def build_index(kb_dir: Path, cache: dict | None = None) -> KbIndex:
    """Extract + chunk every document in kb_dir deterministically."""
    files = _scan_sources(kb_dir)
    chunks: list[KbChunk] = []
    for doc_index, path in enumerate(files):
        text = extract_text(path)
        if not text.strip():
            continue
        chunks.extend(chunk_document(path.name, doc_index, text))
        del text
    gc.collect()
    return KbIndex(
        chunks=chunks,
        sources=_source_records(files, cache),
        built_at=time.time(),
    )


def _current_source_set(
    files: list[Path], cache: dict,
) -> set[tuple[str, str, int]]:
    """Set of (name, content-sha, size) for the current files, using the hash
    cache so unchanged files are not re-hashed."""
    from kb.file_sig import file_sha

    current: set[tuple[str, str, int]] = set()
    for p in files:
        try:
            size = int(p.stat().st_size)
        except OSError:
            size = 0
        current.add((p.name, file_sha(p, cache), size))
    return current


def _index_is_current(index_path: Path, files: list[Path]) -> bool:
    if not index_path.exists():
        return False
    try:
        from kb.kb_crypto import read_decrypted_text
        text = read_decrypted_text(index_path)
        if text is None:
            return False
        data = json.loads(text)
    except Exception:
        return False
    # An index built by an older extractor must be rebuilt even if the source
    # files are byte-for-byte unchanged (so improved extraction takes effect).
    if int(data.get("extractor_version", 1)) != EXTRACTOR_VERSION:
        return False
    # Compare on CONTENT (name, sha, size). An index predating content hashing
    # has empty shas, so it won't match and rebuilds exactly once on upgrade.
    cached = {
        (s.get("name"), str(s.get("sha", "") or ""), int(s.get("size", 0)))
        for s in data.get("sources", [])
    }
    if any(not sha for _n, sha, _sz in cached):
        return False
    from kb.file_sig import (
        hash_cache_path,
        load_hash_cache,
        prune_hash_cache,
        save_hash_cache,
    )

    cache_path = hash_cache_path(index_path)
    cache = load_hash_cache(cache_path)
    current = _current_source_set(files, cache)
    prune_hash_cache(cache, files)
    save_hash_cache(cache_path, cache)
    return cached == current


def _load_index(index_path: Path) -> KbIndex:
    from kb.kb_crypto import read_decrypted_text
    text = read_decrypted_text(index_path)
    if text is None:
        raise OSError("Cannot decrypt kb_index")
    data = json.loads(text)
    chunks = [
        KbChunk(
            chunk_id=str(c.get("chunk_id", "")),
            doc=str(c.get("doc", "")),
            title=str(c.get("title", "")),
            text=str(c.get("text", "")),
            n_chars=int(c.get("n_chars", 0) or 0),
            context=str(c.get("context", "") or ""),
        )
        for c in data.get("chunks", [])
    ]
    sources = [
        KbSource(
            name=str(s.get("name", "")),
            mtime=float(s.get("mtime", 0.0) or 0.0),
            size=int(s.get("size", 0) or 0),
            sha=str(s.get("sha", "") or ""),
        )
        for s in data.get("sources", [])
    ]
    return KbIndex(
        chunks=chunks, sources=sources,
        built_at=float(data.get("built_at", 0.0) or 0.0),
    )


def load_or_build_index(kb_dir: Path, index_path: Path) -> KbIndex:
    """Return a cached index if the documents are unchanged, otherwise
    rebuild and persist. Never raises on cache problems."""
    from kb.file_sig import hash_cache_path, load_hash_cache, save_hash_cache

    files = _scan_sources(kb_dir)
    if _index_is_current(index_path, files):
        try:
            return _load_index(index_path)
        except (OSError, json.JSONDecodeError, ValueError):
            pass
    cache_path = hash_cache_path(index_path)
    cache = load_hash_cache(cache_path)
    index = build_index(kb_dir, cache)
    try:
        from kb.kb_crypto import write_encrypted_text
        index_path.parent.mkdir(parents=True, exist_ok=True)
        write_encrypted_text(
            index_path, json.dumps(index.to_dict(), ensure_ascii=True)
        )
        save_hash_cache(cache_path, cache)
    except OSError:
        pass
    return index
