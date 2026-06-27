"""
defects/parser.py
Programmatic parsing of defect documents (Word, Excel, PowerPoint, PDF)
into structured defect records. Adaptive heading detection handles
variable document formats from different authors.

Pipeline:
    1. Extract text + images from the source document (programmatic)
    2. Attempt structured parsing via heading detection heuristics
    3. If heuristic parsing fails or yields < 1 defect, fall back to LLM
    4. Return a list of ParsedDefect records

Each ParsedDefect contains: parent_id, title, description, repro_steps,
severity, and optional embedded images (as base64).
"""

from __future__ import annotations

import base64
import io
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable

LogFn = Callable[[str], None]

# Common heading patterns people use for defect fields (case-insensitive).
_HEADING_PATTERNS: dict[str, list[str]] = {
    "parent_id": [
        r"parent\s*(?:work\s*)?(?:item)?\s*(?:id|#|number)?",
        r"work\s*item\s*(?:id|#)",
        r"user\s*story\s*(?:id|#)?",
        r"(?:linked?|related)\s*(?:wi|work\s*item)",
    ],
    "title": [
        r"title",
        r"bug\s*(?:title|name|summary)",
        r"defect\s*(?:title|name|summary)",
        r"summary",
        r"issue\s*(?:title|name)?",
    ],
    "description": [
        r"description",
        r"details?",
        r"bug\s*description",
        r"defect\s*description",
        r"issue\s*description",
    ],
    "repro_steps": [
        r"(?:repro(?:duction)?|steps?\s*to\s*(?:repro(?:duce)?|repeat|replicate))",
        r"steps?",
        r"how\s*to\s*repro(?:duce)?",
        r"scenario",
        r"test\s*steps?",
        r"procedure",
    ],
    "severity": [
        r"severity",
        r"priority",
        r"sev(?:erity)?",
        r"impact",
        r"criticality",
    ],
    "expected": [
        r"expected\s*(?:result|behavior|outcome)?",
        r"should\s*(?:be|happen)",
    ],
    "actual": [
        r"actual\s*(?:result|behavior|outcome)?",
        r"what\s*(?:happened|occurs)",
        r"observed",
    ],
}

# Pre-compile heading regexes.
_COMPILED_HEADINGS: dict[str, list[re.Pattern[str]]] = {
    k: [re.compile(p, re.IGNORECASE) for p in patterns]
    for k, patterns in _HEADING_PATTERNS.items()
}


@dataclass(slots=True)
class DefectImage:
    filename: str
    data_b64: str
    mime_type: str = "image/png"


@dataclass(slots=True)
class ParsedDefect:
    parent_id: int = 0
    title: str = ""
    description: str = ""
    repro_steps: str = ""
    severity: str = ""
    expected_result: str = ""
    actual_result: str = ""
    images: list[DefectImage] = field(default_factory=list)

    @property
    def is_valid(self) -> bool:
        return bool(self.title.strip())


def _match_heading(text: str) -> str | None:
    """Match a text string against known heading patterns. Returns the
    field key or None."""
    cleaned = text.strip().rstrip(":").strip()
    if not cleaned or len(cleaned) > 60:
        return None
    for field_key, patterns in _COMPILED_HEADINGS.items():
        for pat in patterns:
            if pat.fullmatch(cleaned):
                return field_key
    return None


def _extract_parent_id(text: str) -> int:
    """Extract a numeric work item ID from text."""
    m = re.search(r"#?(\d{1,9})", text.strip())
    return int(m.group(1)) if m else 0


# -----------------------------------------------------------------
# Document extraction
# -----------------------------------------------------------------
def _extract_word(path: Path) -> tuple[str, list[DefectImage]]:
    """Extract text and images from a Word document."""
    from docx import Document
    doc = Document(str(path))
    parts: list[str] = []
    images: list[DefectImage] = []

    for para in doc.paragraphs:
        style = (para.style.name or "").lower() if para.style else ""
        text = para.text.strip()
        if "heading" in style and text:
            parts.append(f"\n## {text}\n")
        elif text:
            parts.append(text)

    # Extract images from the document's relationships.
    for rel in doc.part.rels.values():
        if "image" in (rel.reltype or ""):
            try:
                blob = rel.target_part.blob
                ext = Path(rel.target_part.partname).suffix.lower()
                mime = f"image/{ext.lstrip('.')}" if ext else "image/png"
                images.append(DefectImage(
                    filename=Path(rel.target_part.partname).name,
                    data_b64=base64.b64encode(blob).decode("ascii"),
                    mime_type=mime,
                ))
            except Exception:
                pass

    return "\n".join(parts), images


def _extract_excel(path: Path) -> tuple[str, list[DefectImage]]:
    """Extract text from an Excel workbook (table-oriented defects)."""
    import openpyxl
    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    parts: list[str] = []
    images: list[DefectImage] = []

    for ws in wb.worksheets:
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        # First row as headers
        headers = [str(c or "").strip() for c in rows[0]]
        parts.append("## " + " | ".join(headers))
        for row in rows[1:]:
            cells = [str(c or "").strip() for c in row]
            if any(cells):
                parts.append(" | ".join(cells))

    wb.close()
    return "\n".join(parts), images


def _extract_powerpoint(path: Path) -> tuple[str, list[DefectImage]]:
    """Extract text and images from a PowerPoint file."""
    from pptx import Presentation
    prs = Presentation(str(path))
    parts: list[str] = []
    images: list[DefectImage] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        parts.append(f"\n## Slide {slide_num}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            if shape.shape_type == 13:  # Picture
                try:
                    blob = shape.image.blob
                    ext = shape.image.ext or "png"
                    images.append(DefectImage(
                        filename=f"slide{slide_num}_{shape.name}.{ext}",
                        data_b64=base64.b64encode(blob).decode("ascii"),
                        mime_type=f"image/{ext}",
                    ))
                except Exception:
                    pass

    return "\n".join(parts), images


def _extract_pdf(path: Path) -> tuple[str, list[DefectImage]]:
    """Extract text and images from a PDF."""
    import pypdf
    reader = pypdf.PdfReader(str(path))
    parts: list[str] = []
    images: list[DefectImage] = []

    for page_num, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"--- Page {page_num} ---\n{text}")
        # Extract images (per-image error handling).
        try:
            page_images = page.images or []
        except Exception:
            page_images = []
        for img_key in page_images:
            try:
                blob = img_key.data
                images.append(DefectImage(
                    filename=f"page{page_num}_{img_key.name}",
                    data_b64=base64.b64encode(blob).decode("ascii"),
                    mime_type="image/png",
                ))
            except Exception:
                pass

    return "\n".join(parts), images


def extract_document(path: Path, on_log: LogFn | None = None) -> tuple[str, list[DefectImage]]:
    """Route extraction to the correct handler based on file extension."""
    ext = path.suffix.lower()
    _log = on_log or (lambda _: None)
    try:
        if ext == ".docx":
            _log(f"[INFO] Parsing Word document: {path.name}")
            return _extract_word(path)
        elif ext == ".xlsx":
            _log(f"[INFO] Parsing Excel workbook: {path.name}")
            return _extract_excel(path)
        elif ext == ".pptx":
            _log(f"[INFO] Parsing PowerPoint: {path.name}")
            return _extract_powerpoint(path)
        elif ext in (".doc", ".xls", ".ppt"):
            _log(f"[INFO] Parsing legacy {ext} via fallback: {path.name}")
            return _extract_legacy(path)
        elif ext == ".pdf":
            _log(f"[INFO] Parsing PDF: {path.name}")
            return _extract_pdf(path)
        else:
            _log(f"[WARN] Unsupported file type: {ext}")
            return "", []
    except Exception as e:
        _log(f"[ERROR] Extraction failed for {path.name}: {e!r}")
        return "", []


def _extract_legacy(path: Path) -> tuple[str, list[DefectImage]]:
    """Extract text from legacy Office formats (.doc, .xls, .ppt) using
    the KB legacy extractor. No image extraction for legacy formats."""
    from kb.legacy_docs import extract_legacy_text
    text = extract_legacy_text(path)
    return text, []


# -----------------------------------------------------------------
# Heuristic parsing
# -----------------------------------------------------------------
def _parse_structured_text(text: str) -> list[ParsedDefect]:
    """Attempt to parse extracted text into defects using heading detection."""
    defects: list[ParsedDefect] = []
    lines = text.splitlines()
    current: ParsedDefect | None = None
    current_field: str = ""
    field_buffer: list[str] = []

    def _flush_field() -> None:
        nonlocal current, current_field, field_buffer
        if current is None or not current_field or not field_buffer:
            field_buffer = []
            return
        value = "\n".join(field_buffer).strip()
        if current_field == "parent_id":
            current.parent_id = _extract_parent_id(value)
        elif current_field == "title":
            current.title = value
        elif current_field == "description":
            current.description = value
        elif current_field == "repro_steps":
            current.repro_steps = value
        elif current_field == "severity":
            current.severity = value
        elif current_field == "expected":
            current.expected_result = value
        elif current_field == "actual":
            current.actual_result = value
        field_buffer = []

    # Detect if this is a table format (Excel-style with | separators).
    pipe_lines = [l for l in lines if "|" in l and l.count("|") >= 2]
    if len(pipe_lines) > 2:
        return _parse_table_format(lines)

    for line in lines:
        stripped = line.strip()
        # Detect heading-style lines (## heading, **heading**, or plain heading:)
        heading_text = ""
        if stripped.startswith("##"):
            heading_text = stripped.lstrip("#").strip()
        elif stripped.endswith(":") and len(stripped) < 50:
            heading_text = stripped
        elif re.match(r"^\*\*(.+?)\*\*\s*:?$", stripped):
            heading_text = re.sub(r"\*\*", "", stripped).strip(": ")

        if heading_text:
            field_match = _match_heading(heading_text)
            if field_match == "title" and current is not None:
                # New defect boundary when we see a new title heading.
                _flush_field()
                if current.is_valid:
                    defects.append(current)
                current = ParsedDefect()
                current_field = "title"
                continue
            elif field_match:
                _flush_field()
                if current is None:
                    current = ParsedDefect()
                current_field = field_match
                continue

        # Check for defect delimiter patterns (numbered defects, separators).
        defect_num = re.match(
            r"^(?:defect|bug|issue)\s*#?\s*(\d+)", stripped, re.IGNORECASE
        )
        if defect_num or stripped.startswith("==="):
            _flush_field()
            if current is not None and current.is_valid:
                defects.append(current)
            current = ParsedDefect()
            current_field = ""
            continue

        # Accumulate content for the current field.
        if current_field:
            field_buffer.append(stripped)
        elif stripped and current is None:
            # First content before any heading - could be a title.
            current = ParsedDefect()
            current_field = "title"
            field_buffer.append(stripped)

    _flush_field()
    if current is not None and current.is_valid:
        defects.append(current)

    return defects


def _parse_table_format(lines: list[str]) -> list[ParsedDefect]:
    """Parse a pipe-delimited table format (from Excel extraction)."""
    defects: list[ParsedDefect] = []
    # Find header row.
    header_idx = -1
    headers: list[str] = []
    for i, line in enumerate(lines):
        if "|" in line and line.count("|") >= 2:
            cells = [c.strip() for c in line.split("|")]
            # Check if this looks like a header row.
            matched = sum(1 for c in cells if _match_heading(c) is not None)
            if matched >= 2:
                header_idx = i
                headers = cells
                break

    if header_idx < 0 or not headers:
        return []

    # Map header positions to field keys.
    col_map: dict[int, str] = {}
    for i, h in enumerate(headers):
        key = _match_heading(h)
        if key:
            col_map[i] = key

    # Parse data rows.
    for line in lines[header_idx + 1:]:
        if "|" not in line:
            continue
        cells = [c.strip() for c in line.split("|")]
        if not any(cells):
            continue
        defect = ParsedDefect()
        for col_idx, field_key in col_map.items():
            if col_idx >= len(cells):
                continue
            value = cells[col_idx]
            if field_key == "parent_id":
                defect.parent_id = _extract_parent_id(value)
            elif field_key == "title":
                defect.title = value
            elif field_key == "description":
                defect.description = value
            elif field_key == "repro_steps":
                defect.repro_steps = value
            elif field_key == "severity":
                defect.severity = value
            elif field_key == "expected":
                defect.expected_result = value
            elif field_key == "actual":
                defect.actual_result = value
        if defect.is_valid:
            defects.append(defect)

    return defects


# -----------------------------------------------------------------
# LLM fallback parsing
# -----------------------------------------------------------------
_LLM_PARSE_SYSTEM = (
    "You are a structured data extractor. Given a document containing "
    "software defect/bug reports, extract each defect into a JSON array. "
    "Each defect object must have these fields:\n"
    "- parent_id: integer (the parent work item ID, 0 if not found)\n"
    "- title: string (concise bug title)\n"
    "- description: string (detailed description)\n"
    "- repro_steps: string (numbered steps to reproduce)\n"
    "- severity: string (Critical/High/Medium/Low, infer if not stated)\n"
    "- expected_result: string (what should happen)\n"
    "- actual_result: string (what actually happens)\n\n"
    "Output ONLY a JSON array of defect objects. No prose, no code fence."
)


async def parse_with_llm_async(
    client: Any,
    model: str,
    text: str,
    on_log: LogFn | None = None,
) -> list[ParsedDefect]:
    """Fall back to LLM parsing when programmatic parsing fails."""
    import json
    _log = on_log or (lambda _: None)
    _log("[INFO] Programmatic parsing insufficient; using LLM to extract "
         "defect structure...")

    truncated = text[:80000]
    try:
        result = await client.complete_async(
            model=model, system=_LLM_PARSE_SYSTEM,
            user=f"Extract all defects from the following document:\n\n"
                 f"{truncated}",
            max_tokens=8000, temperature=0.0,
        )
    except Exception as e:
        _log(f"[ERROR] LLM parsing failed: {e!r}")
        return []

    raw = (result.text or "").strip()
    # Try to extract JSON array.
    start = raw.find("[")
    end = raw.rfind("]")
    if start == -1 or end == -1:
        _log("[ERROR] LLM did not return a valid JSON array.")
        return []

    try:
        items = json.loads(raw[start:end + 1])
    except json.JSONDecodeError as e:
        _log(f"[ERROR] LLM JSON parse error: {e}")
        return []

    defects: list[ParsedDefect] = []
    for item in items:
        if not isinstance(item, dict):
            continue
        raw_pid = item.get("parent_id", 0)
        try:
            pid = int(float(raw_pid)) if raw_pid else 0
        except (ValueError, TypeError):
            pid = 0
        defects.append(ParsedDefect(
            parent_id=pid,
            title=str(item.get("title", "") or ""),
            description=str(item.get("description", "") or ""),
            repro_steps=str(item.get("repro_steps", "") or ""),
            severity=str(item.get("severity", "") or ""),
            expected_result=str(item.get("expected_result", "") or ""),
            actual_result=str(item.get("actual_result", "") or ""),
        ))
    valid = [d for d in defects if d.is_valid]
    _log(f"[SUCCESS] LLM extracted {len(valid)} defect(s).")
    return valid


# -----------------------------------------------------------------
# Public entry point
# -----------------------------------------------------------------
async def parse_defect_documents_async(
    paths: list[Path],
    on_log: LogFn | None = None,
    on_progress: Callable[[int, int], None] | None = None,
    llm_client: Any | None = None,
    llm_model: str = "",
) -> list[ParsedDefect]:
    """Parse one or more documents into defect records. Falls back to LLM
    when programmatic parsing yields no results."""
    _log = on_log or (lambda _: None)
    all_defects: list[ParsedDefect] = []
    all_images: list[DefectImage] = []

    for i, path in enumerate(paths):
        if on_progress:
            on_progress(i, len(paths))
        text, images = extract_document(path, on_log=on_log)
        all_images.extend(images)

        if not text.strip():
            _log(f"[WARN] No text extracted from {path.name}; skipping.")
            continue

        # Attempt programmatic parsing first.
        defects = _parse_structured_text(text)
        if defects:
            _log(f"[SUCCESS] Parsed {len(defects)} defect(s) from "
                 f"{path.name} (programmatic).")
            all_defects.extend(defects)
        elif llm_client is not None and llm_model:
            # LLM fallback.
            defects = await parse_with_llm_async(
                llm_client, llm_model, text, on_log=on_log
            )
            all_defects.extend(defects)
        else:
            _log(f"[WARN] Could not parse defects from {path.name} and "
                 f"no LLM available for fallback.")

    if on_progress:
        on_progress(len(paths), len(paths))

    # Distribute images across defects (best-effort: round-robin, capped).
    if all_images and all_defects:
        n_img = len(all_images)
        n_def = len(all_defects)
        per_defect = max(1, n_img // n_def)
        idx = 0
        for i, defect in enumerate(all_defects):
            if idx >= n_img:
                break
            end = n_img if i == n_def - 1 else min(idx + per_defect, n_img)
            defect.images = all_images[idx:end]
            idx = end

    _log(f"[INFO] Total: {len(all_defects)} defect(s) parsed from "
         f"{len(paths)} file(s).")
    return all_defects
